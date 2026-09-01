import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette (Modern Premium Slate & Cyan Theme)
COLOR_BG = RGBColor(15, 23, 42)          # Dark Navy #0F172A
COLOR_CARD = RGBColor(30, 41, 59)        # Slate 800 #1E293B
COLOR_CARD_BORDER = RGBColor(51, 65, 85) # Slate 700
COLOR_PRIMARY = RGBColor(0, 229, 255)     # Cyan #00E5FF
COLOR_ACCENT = RGBColor(0, 230, 118)     # Neon Green #00E676
COLOR_WARN = RGBColor(255, 171, 0)       # Amber #FFAB00
COLOR_TEXT_MAIN = RGBColor(248, 250, 252)# White/Light
COLOR_TEXT_SUB = RGBColor(148, 163, 184) # Slate 400

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

def add_header(slide, title_text, category_text="대시보드 배포 가이드"):
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8), Inches(0.4))
    tf_c = cat_box.text_frame
    tf_c.word_wrap = True
    p_c = tf_c.paragraphs[0]
    p_c.text = f"■ {category_text.upper()}"
    p_c.font.size = Pt(11)
    p_c.font.bold = True
    p_c.font.color.rgb = COLOR_PRIMARY
    p_c.font.name = "맑은 고딕"

    t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.8))
    tf_t = t_box.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    p_t.text = title_text
    p_t.font.size = Pt(22)
    p_t.font.bold = True
    p_t.font.color.rgb = COLOR_TEXT_MAIN
    p_t.font.name = "맑은 고딕"

blank_layout = prs.slide_layouts[6]

# ==========================================
# SLIDE 1: 표지 (Cover)
# ==========================================
s1 = prs.slides.add_slide(blank_layout)
set_slide_background(s1)

tbox = s1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(3.0))
tf = tbox.text_frame
tf.word_wrap = True

p0 = tf.paragraphs[0]
p0.text = "📊 기술본부 업무량 및 지원 시간 분석 대시보드"
p0.font.size = Pt(26)
p0.font.bold = True
p0.font.color.rgb = COLOR_PRIMARY
p0.font.name = "맑은 고딕"
p0.space_after = Pt(14)

p1 = tf.add_paragraph()
p1.text = "다른 PC 간편 설치 및 실행 / 업데이트 가이드"
p1.font.size = Pt(36)
p1.font.bold = True
p1.font.color.rgb = COLOR_TEXT_MAIN
p1.font.name = "맑은 고딕"
p1.space_after = Pt(20)

p2 = tf.add_paragraph()
p2.text = "GitHub 저장소(PublicRepo) 연동 및 원클릭 배치 파일(.bat) 구동 매뉴얼"
p2.font.size = Pt(16)
p2.font.color.rgb = COLOR_TEXT_SUB
p2.font.name = "맑은 고딕"

foot = s1.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(11.3), Inches(0.8))
tf_f = foot.text_frame
pf = tf_f.paragraphs[0]
pf.text = "기술본부 | 독립 로컬 SQLite 및 PC 카카오톡 10분 자동 수집 시스템 탑재"
pf.font.size = Pt(13)
pf.font.color.rgb = COLOR_ACCENT
pf.font.name = "맑은 고딕"

# ==========================================
# SLIDE 2: 1단계 - 다른 PC 최초 1회 사전 준비
# ==========================================
s2 = prs.slides.add_slide(blank_layout)
set_slide_background(s2)
add_header(s2, "Step 1. 다른 PC 최초 1회 사전 준비")

cards_data = [
    ("1. Python 3.10+ 설치", [
        "• Python 공식 홈페이지(python.org)에서 다운로드",
        "• ★ 설치 화면에서 [Add python.exe to PATH] 반드시 체크!",
        "• 설치 완료 후 터미널에서 python --version 확인"
    ], COLOR_PRIMARY),
    ("2. 저장소(PublicRepo) 복사", [
        "• 방법 A (Git 사용): git clone https://github.com/newprim82/PublicRepo.git",
        "• 방법 B (ZIP 다운로드): GitHub 페이지에서 [Code] > [Download ZIP] 후 압축 해제",
        "• 원하는 폴더(예: C:\\Python)에 배치"
    ], COLOR_ACCENT),
    ("3. 카카오톡 대화방 열기", [
        "• 해당 PC에서 PC 카카오톡 로그인",
        "• [기술본부] 업무공유방 창을 더블클릭하여 열어둠",
        "• 대시보드가 10분마다 자동으로 대화를 긁어와 동기화합니다."
    ], COLOR_WARN)
]

for i, (title, bullets, accent) in enumerate(cards_data):
    left = Inches(0.8 + i * 3.95)
    top = Inches(1.8)
    width = Inches(3.8)
    height = Inches(4.8)
    
    shape = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_CARD
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.5)
    
    tb = s2.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = accent
    p.font.name = "맑은 고딕"
    p.space_after = Pt(14)
    
    for b in bullets:
        pb = tf.add_paragraph()
        pb.text = b
        pb.font.size = Pt(13)
        pb.font.color.rgb = COLOR_TEXT_MAIN if ("★" in b or "방법" in b) else COLOR_TEXT_SUB
        pb.font.name = "맑은 고딕"
        pb.space_after = Pt(8)

# ==========================================
# SLIDE 3: 2단계 - 최초 실행 (setup_and_run.bat)
# ==========================================
s3 = prs.slides.add_slide(blank_layout)
set_slide_background(s3)
add_header(s3, "Step 2. 최초 1회 실행 방법 (원클릭 자동 설정)")

left_box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.7), Inches(4.8))
left_box.fill.solid()
left_box.fill.fore_color.rgb = COLOR_CARD
left_box.line.color.rgb = COLOR_ACCENT
left_box.line.width = Pt(2)

tb_l = s3.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.3), Inches(4.4))
tf_l = tb_l.text_frame
tf_l.word_wrap = True

p = tf_l.paragraphs[0]
p.text = "🚀 setup_and_run.bat 더블 클릭!"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = COLOR_ACCENT
p.font.name = "맑은 고딕"
p.space_after = Pt(16)

bullets_l = [
    "1. 폴더 안에 있는 setup_and_run.bat 파일을 더블 클릭합니다.",
    "2. 대시보드 구동에 필요한 필수 패키지가 100% 자동 설치됩니다:",
    "   - streamlit (웹 UI 및 대시보드 프레임워크)",
    "   - pandas, plotly (데이터 통계 집계 및 인터랙티브 차트)",
    "   - pywin32, uiautomation (카카오톡 자동 수집기)",
    "3. 설치 완료 후 웹 브라우저(http://localhost:8501)가 자동으로 열리며 대시보드가 즉시 시작됩니다!"
]
for b in bullets_l:
    pb = tf_l.add_paragraph()
    pb.text = b
    pb.font.size = Pt(13.5)
    pb.font.color.rgb = COLOR_TEXT_MAIN
    pb.font.name = "맑은 고딕"
    pb.space_after = Pt(10)

right_box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
right_box.fill.solid()
right_box.fill.fore_color.rgb = COLOR_CARD
right_box.line.color.rgb = COLOR_PRIMARY
right_box.line.width = Pt(1.5)

tb_r = s3.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.4))
tf_r = tb_r.text_frame
tf_r.word_wrap = True

p_r = tf_r.paragraphs[0]
p_r.text = "💾 기존 데이터(2,081건) 즉시 연동"
p_r.font.size = Pt(20)
p_r.font.bold = True
p_r.font.color.rgb = COLOR_PRIMARY
p_r.font.name = "맑은 고딕"
p_r.space_after = Pt(16)

bullets_r = [
    "• GitHub에 데이터베이스(data/worklog.db)가 함께 포함되어 있어,",
    "  별도의 DB 세팅 없이 기존 작업 기록 2,081건이 즉시 로드됩니다.",
    "• 팀원별 소속팀(기술 1팀~3팀, PI팀) 및 직급(사원~수석) 정보도 완벽하게 유지됩니다.",
    "• 과중 근무(주 40h/52h) 알림 배너 및 보상 휴가 기록도 바로 확인 가능합니다."
]
for b in bullets_r:
    pb = tf_r.add_paragraph()
    pb.text = b
    pb.font.size = Pt(13.5)
    pb.font.color.rgb = COLOR_TEXT_MAIN
    pb.font.name = "맑은 고딕"
    pb.space_after = Pt(10)

# ==========================================
# SLIDE 4: 3단계 - 일상 실행 및 최신 업데이트
# ==========================================
s4 = prs.slides.add_slide(blank_layout)
set_slide_background(s4)
add_header(s4, "Step 3. 평소 실행 및 최신 업데이트 방법")

b1 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.7), Inches(4.8))
b1.fill.solid()
b1.fill.fore_color.rgb = COLOR_CARD
b1.line.color.rgb = COLOR_PRIMARY
b1.line.width = Pt(2)

tb1 = s4.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.3), Inches(4.4))
tf1 = tb1.text_frame
tf1.word_wrap = True

p1 = tf1.paragraphs[0]
p1.text = "⚡ 평소 일상 실행 (매일 사용 시)"
p1.font.size = Pt(19)
p1.font.bold = True
p1.font.color.rgb = COLOR_PRIMARY
p1.font.name = "맑은 고딕"
p1.space_after = Pt(14)

b_list1 = [
    "▶ 실행 파일: run_dashboard.bat 더블 클릭!",
    "• 추가 설치 과정 없이 1초 만에 대시보드가 열립니다.",
    "• 백그라운드에서 10분마다 [기술본부] 업무공유방 대화를 자동 수집합니다.",
    "• 메인 상단 또는 사이드바의 [⚡ 즉시 수집] 버튼을 누르면 지금 즉시 최신 대화를 긁어옵니다."
]
for b in b_list1:
    pb = tf1.add_paragraph()
    pb.text = b
    pb.font.size = Pt(13.5)
    pb.font.color.rgb = COLOR_TEXT_MAIN if "▶" in b else COLOR_TEXT_SUB
    pb.font.name = "맑은 고딕"
    pb.space_after = Pt(8)

b2 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
b2.fill.solid()
b2.fill.fore_color.rgb = COLOR_CARD
b2.line.color.rgb = COLOR_ACCENT
b2.line.width = Pt(2)

tb2 = s4.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.4))
tf2 = tb2.text_frame
tf2.word_wrap = True

p2 = tf2.paragraphs[0]
p2.text = "🔄 내 PC에서 수정한 기능 반영 (업데이트 시)"
p2.font.size = Pt(19)
p2.font.bold = True
p2.font.color.rgb = COLOR_ACCENT
p2.font.name = "맑은 고딕"
p2.space_after = Pt(14)

b_list2 = [
    "▶ 실행 파일: update_and_run.bat 더블 클릭!",
    "• 내 PC에서 기능을 추가/수정하고 GitHub에 올렸을 때 사용합니다.",
    "• GitHub에서 최신 소스코드를 자동으로 git pull 받아옵니다.",
    "• 신규 패키지 설치 및 대시보드 실행까지 한 번에 자동 진행됩니다!"
]
for b in b_list2:
    pb = tf2.add_paragraph()
    pb.text = b
    pb.font.size = Pt(13.5)
    pb.font.color.rgb = COLOR_TEXT_MAIN if "▶" in b else COLOR_TEXT_SUB
    pb.font.name = "맑은 고딕"
    pb.space_after = Pt(8)

# ==========================================
# SLIDE 5: 핵심 기능 및 팁 요약
# ==========================================
s5 = prs.slides.add_slide(blank_layout)
set_slide_background(s5)
add_header(s5, "💡 핵심 기능 및 편리한 사용 꿀팁", "핵심 가이드")

features = [
    ("🤖 카카오톡 10분 증분 자동 수집", "대화방 창을 열어두면 10분마다 안전하게(100% 읽기 전용) 긁어옵니다. 채팅방에 아무런 글자도 쳐지지 않습니다.", COLOR_PRIMARY),
    ("⚡ 원클릭 즉시 수집", "10분을 기다리지 않고 지금 올라온 메시지를 즉시 반영하려면 [⚡ 즉시 수집] 버튼을 클릭하세요.", COLOR_ACCENT),
    ("🚨 과중 근무 실시간 감지 & 보상 휴가", "주 40h(주황색) / 주 52h(빨간색) 초과 인원이 반짝이는 네온 배너에 자동 감지되며, 클릭 시 보상 휴가 모달이 뜹니다.", COLOR_WARN),
    ("👔 직급 및 팀원 관리", "사원 / 대리 / 과장 / 수석 회사 표준 4대 직급이 적용되어 있으며, 필터와 팝업에서 직급별/팀별 분석이 가능합니다.", COLOR_PRIMARY)
]

for i, (f_title, f_desc, f_col) in enumerate(features):
    row = i // 2
    col = i % 2
    left = Inches(0.8 + col * 5.95)
    top = Inches(1.8 + row * 2.5)
    width = Inches(5.8)
    height = Inches(2.2)
    
    shape = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_CARD
    shape.line.color.rgb = f_col
    shape.line.width = Pt(1.5)
    
    tb = s5.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2), width - Inches(0.5), height - Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = f_title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = f_col
    p.font.name = "맑은 고딕"
    p.space_after = Pt(8)
    
    p_d = tf.add_paragraph()
    p_d.text = f_desc
    p_d.font.size = Pt(13)
    p_d.font.color.rgb = COLOR_TEXT_MAIN
    p_d.font.name = "맑은 고딕"

output_path = r"c:\Python\work-time-dashboard\대시보드_다른PC_설치및실행_가이드.pptx"
prs.save(output_path)
print(f"[PPT 생성 성공] {output_path}")
